"""
Agentic - Document Reader Skills
=================================
Provides skills for reading binary document formats that the plain
``read_file`` skill cannot handle:

  • read_pdf   – PDF files via pypdf
  • read_excel – Excel workbooks (.xlsx / .xls) via openpyxl
  • read_word  – Word documents (.docx) via python-docx
  • read_pptx  – PowerPoint presentations (.pptx) via python-pptx

All skills share the same path-safety rules as the filesystem skill
(blocked system directories, symlink-resolved paths).

Vision / image reading is intentionally not included here.
When a vision-capable model is loaded, image inputs can be wired
directly into ModelNexus.stream() as a future enhancement.

Dependencies (install once):
    pip install pypdf openpyxl python-docx python-pptx
"""
from __future__ import annotations

import os
from pathlib import Path
from typing import Any

from skills.base import SkillBase
from utils.logger import build_logger

log = build_logger("agentic.skill.doc_reader")

# ── path safety (mirrors filesystem.py) ─────────────────────────────────────

_BLOCKED_DIRS: frozenset[Path] = frozenset({
    Path("/etc"),
    Path("/bin"),
    Path("/sbin"),
    Path("/usr/bin"),
    Path("/usr/sbin"),
    Path("/sys"),
    Path("/proc"),
    Path("/dev"),
    Path("/boot"),
    Path("/root"),
})


def _safe_path(raw: str) -> Path:
    p = Path(raw).expanduser().resolve()
    ancestors = [p, *p.parents]
    for ancestor in ancestors:
        if ancestor in _BLOCKED_DIRS:
            raise PermissionError(
                f"Access to '{p}' is not allowed (falls under blocked directory '{ancestor}')"
            )
    return p


def _check_exists(p: Path, suffix_hint: str) -> None:
    if not p.exists():
        raise FileNotFoundError(f"File not found: {p}")
    if not p.is_file():
        raise ValueError(f"Path is not a file: {p}")


# ── PDF ──────────────────────────────────────────────────────────────────────

class ReadPdfSkill(SkillBase):
    name = "read_pdf"
    description = (
        "Extract and return the text content of a PDF file. "
        "Requires the 'pypdf' package (pip install pypdf)."
    )
    parameters = {
        "path": {"type": "string", "description": "Absolute or ~ path to the PDF file"},
        "max_chars": {
            "type": "integer",
            "description": "Maximum characters to return (default 8000)",
        },
        "pages": {
            "type": "string",
            "description": (
                "Comma-separated 1-based page numbers or 'all' (default 'all'). "
                "Example: '1,2,5'"
            ),
        },
    }
    required = ["path"]
    tags = ["document", "filesystem"]

    async def execute(
        self, path: str, max_chars: int = 8000, pages: str = "all"
    ) -> str:
        try:
            from pypdf import PdfReader
        except ImportError:
            raise RuntimeError(
                "pypdf is required for PDF reading. Run: pip install pypdf"
            )

        p = _safe_path(path)
        _check_exists(p, ".pdf")

        reader = PdfReader(str(p))
        total_pages = len(reader.pages)

        if pages == "all":
            page_indices = list(range(total_pages))
        else:
            page_indices = []
            for token in pages.split(","):
                token = token.strip()
                if token.isdigit():
                    idx = int(token) - 1
                    if 0 <= idx < total_pages:
                        page_indices.append(idx)

        parts: list[str] = []
        for idx in page_indices:
            text = reader.pages[idx].extract_text() or ""
            if text:
                parts.append(f"--- Page {idx + 1} ---\n{text.strip()}")

        result = "\n\n".join(parts) if parts else "(No text could be extracted from this PDF)"
        if len(result) > max_chars:
            result = result[:max_chars] + f"\n... [truncated at {max_chars} chars]"

        log.info("read_pdf: %s (%d pages read)", p.name, len(page_indices))
        return result


# ── Excel ────────────────────────────────────────────────────────────────────

class ReadExcelSkill(SkillBase):
    name = "read_excel"
    description = (
        "Read an Excel workbook (.xlsx) and return its content as "
        "Markdown tables. Requires 'openpyxl' (pip install openpyxl)."
    )
    parameters = {
        "path": {"type": "string", "description": "Absolute or ~ path to the .xlsx file"},
        "sheet": {
            "type": "string",
            "description": "Sheet name or 1-based sheet index to read (default: all sheets)",
        },
        "max_rows": {
            "type": "integer",
            "description": "Max rows per sheet to return (default 200)",
        },
        "max_chars": {
            "type": "integer",
            "description": "Maximum characters to return in total (default 8000)",
        },
    }
    required = ["path"]
    tags = ["document", "filesystem"]

    async def execute(
        self,
        path: str,
        sheet: str = "",
        max_rows: int = 200,
        max_chars: int = 8000,
    ) -> str:
        try:
            import openpyxl
        except ImportError:
            raise RuntimeError(
                "openpyxl is required for Excel reading. Run: pip install openpyxl"
            )

        p = _safe_path(path)
        _check_exists(p, ".xlsx")

        wb = openpyxl.load_workbook(str(p), read_only=True, data_only=True)

        # Determine which sheets to read
        if sheet:
            if sheet.isdigit():
                idx = int(sheet) - 1
                if idx < 0 or idx >= len(wb.sheetnames):
                    raise ValueError(
                        f"Sheet index {int(sheet)} out of range (workbook has {len(wb.sheetnames)} sheets)"
                    )
                sheet_names = [wb.sheetnames[idx]]
            elif sheet in wb.sheetnames:
                sheet_names = [sheet]
            else:
                raise ValueError(
                    f"Sheet '{sheet}' not found. Available: {', '.join(wb.sheetnames)}"
                )
        else:
            sheet_names = wb.sheetnames

        parts: list[str] = []
        for name in sheet_names:
            ws = wb[name]
            rows_data: list[list[str]] = []
            for row in ws.iter_rows(max_row=max_rows, values_only=True):
                rows_data.append([str(cell) if cell is not None else "" for cell in row])

            if not rows_data:
                parts.append(f"### Sheet: {name}\n_(empty)_")
                continue

            # Build a Markdown table
            header = rows_data[0]
            sep = ["---"] * len(header)
            md_rows = [
                "| " + " | ".join(header) + " |",
                "| " + " | ".join(sep) + " |",
            ]
            for row in rows_data[1:]:
                md_rows.append("| " + " | ".join(row) + " |")

            parts.append(f"### Sheet: {name}\n" + "\n".join(md_rows))

        wb.close()

        result = "\n\n".join(parts)
        if len(result) > max_chars:
            result = result[:max_chars] + f"\n... [truncated at {max_chars} chars]"

        log.info("read_excel: %s (%d sheets)", p.name, len(sheet_names))
        return result


# ── Word (.docx) ─────────────────────────────────────────────────────────────

class ReadWordSkill(SkillBase):
    name = "read_word"
    description = (
        "Extract and return the text content of a Word document (.docx). "
        "Requires 'python-docx' (pip install python-docx)."
    )
    parameters = {
        "path": {"type": "string", "description": "Absolute or ~ path to the .docx file"},
        "max_chars": {
            "type": "integer",
            "description": "Maximum characters to return (default 8000)",
        },
        "include_tables": {
            "type": "boolean",
            "description": "Include table content as Markdown (default true)",
        },
    }
    required = ["path"]
    tags = ["document", "filesystem"]

    async def execute(
        self, path: str, max_chars: int = 8000, include_tables: bool = True
    ) -> str:
        try:
            from docx import Document
        except ImportError:
            raise RuntimeError(
                "python-docx is required for Word reading. Run: pip install python-docx"
            )

        p = _safe_path(path)
        _check_exists(p, ".docx")

        doc = Document(str(p))
        parts: list[str] = []

        for block in doc.element.body:
            tag = block.tag.split("}")[-1] if "}" in block.tag else block.tag

            if tag == "p":
                # Paragraph
                text = "".join(node.text or "" for node in block.iter() if node.text)
                if text.strip():
                    parts.append(text.strip())

            elif tag == "tbl" and include_tables:
                # Table – convert to Markdown
                from docx.oxml.ns import qn
                rows_md: list[str] = []
                for i, tr in enumerate(block.findall(f".//{qn('w:tr')}")):
                    cells = [
                        "".join(n.text or "" for n in tc.iter() if n.text)
                        for tc in tr.findall(f".//{qn('w:tc')}")
                    ]
                    rows_md.append("| " + " | ".join(cells) + " |")
                    if i == 0:
                        rows_md.append("| " + " | ".join(["---"] * len(cells)) + " |")
                if rows_md:
                    parts.append("\n".join(rows_md))

        result = "\n\n".join(parts) if parts else "(No text could be extracted from this document)"
        if len(result) > max_chars:
            result = result[:max_chars] + f"\n... [truncated at {max_chars} chars]"

        log.info("read_word: %s", p.name)
        return result


# ── PowerPoint (.pptx) ───────────────────────────────────────────────────────

class ReadPptxSkill(SkillBase):
    name = "read_pptx"
    description = (
        "Extract and return the text content of a PowerPoint presentation (.pptx). "
        "Requires 'python-pptx' (pip install python-pptx)."
    )
    parameters = {
        "path": {"type": "string", "description": "Absolute or ~ path to the .pptx file"},
        "max_chars": {
            "type": "integer",
            "description": "Maximum characters to return (default 8000)",
        },
        "slides": {
            "type": "string",
            "description": (
                "Comma-separated 1-based slide numbers or 'all' (default 'all'). "
                "Example: '1,3,5'"
            ),
        },
    }
    required = ["path"]
    tags = ["document", "filesystem"]

    async def execute(
        self, path: str, max_chars: int = 8000, slides: str = "all"
    ) -> str:
        try:
            from pptx import Presentation
            from pptx.util import Pt
        except ImportError:
            raise RuntimeError(
                "python-pptx is required for PowerPoint reading. Run: pip install python-pptx"
            )

        p = _safe_path(path)
        _check_exists(p, ".pptx")

        prs = Presentation(str(p))
        total = len(prs.slides)

        if slides == "all":
            slide_indices = list(range(total))
        else:
            slide_indices = []
            for token in slides.split(","):
                token = token.strip()
                if token.isdigit():
                    idx = int(token) - 1
                    if 0 <= idx < total:
                        slide_indices.append(idx)

        parts: list[str] = []
        for idx in slide_indices:
            slide = prs.slides[idx]
            slide_texts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                parts.append(f"--- Slide {idx + 1} ---\n" + "\n".join(slide_texts))

        result = "\n\n".join(parts) if parts else "(No text could be extracted from this presentation)"
        if len(result) > max_chars:
            result = result[:max_chars] + f"\n... [truncated at {max_chars} chars]"

        log.info("read_pptx: %s (%d slides read)", p.name, len(slide_indices))
        return result


# ── Registration ─────────────────────────────────────────────────────────────

def register_all() -> None:
    ReadPdfSkill.register()
    ReadExcelSkill.register()
    ReadWordSkill.register()
    ReadPptxSkill.register()
