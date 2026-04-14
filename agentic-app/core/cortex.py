"""
Agentic - Cortex
================
The central reasoning unit of the Reactive Cortex Architecture (RCA).

The Cortex drives the *Deliberation Pulse* using a ReAct loop
(Reason → Act → Observe → Reason … → Final Answer):

  1. Receive user input via the Signal Lattice.
  2. Assemble context from the Memory Lattice (query-relevance ranked).
  3. Build system prompt and initial message list.
  4. Resolve tool schema (native calling for supported models).
  5. Loop up to react_max_iterations times:
       a. Stream a model response.
       b. Extract any skill-call markers from the response.
       c. If none → Final Answer.  Break.
       d. Run skill invocations in parallel via TaskFibers (Act).
       e. Build an Observation message from results (Observe).
       f. Append assistant turn + observation to message history, loop.
  6. Write final answer to the Memory Lattice.
  7. Emit DELIBERATION_END signal with the full response.

Thread model:
  The Cortex runs its async event loop in a dedicated background thread.
  UI interaction happens via Signal Lattice callbacks only – the Cortex
  never touches tkinter widgets directly.
"""
from __future__ import annotations

import asyncio
import re
import threading
from pathlib import Path
from typing import Callable

from core.task_orchestrator import build_plan, format_plan_block, quality_gate
from core.memory_lattice import MemoryLattice
from core.signal_lattice import SigKind, Signal, lattice
from core.skill_registry import SkillRegistry
from core.task_fabric import FiberStatus, TaskFabric, TaskFiber
from model.gemma_nexus import GemmaNexus, get_assistant_role
from model.prompt_weaver import PromptWeaver, SkillInvocation
from utils.config import cfg
from utils.logger import build_logger

log = build_logger("agentic.cortex")

# Status message template emitted to the UI after each ReAct tool round-trip.
# Shows the current iteration and the names of skills that were executed.
_REACT_ITERATION_STATUS = (
    "\n\n[⚙ Skills executed: {skill_names} — reasoning continues…]\n\n"
)

# Message injected into the conversation when react_max_iterations is exhausted
# but the model is still calling tools (no clean final answer was produced).
# This Reflexion step (Shinn et al., 2023) forces one final synthesis pass with
# tools disabled so the model generates a direct answer from accumulated evidence
# instead of looping further or returning a response that is full of tool calls.
_REFLEXION_PROMPT = (
    "You have reached the maximum number of reasoning steps and cannot call "
    "any more tools. Based on everything you have gathered so far, please "
    "synthesize a clear, direct final answer for the user. "
    "Do not invoke any tools or skills in your response."
)

# Status token emitted to the streaming UI when Reflexion is triggered.
_REFLEXION_STATUS = "\n\n[↺ Reflexion: synthesizing findings into a final answer…]\n\n"

# Fraction of the context limit at which we warn and trim the oldest messages.
# 0.85 leaves ~15 % headroom for the model's reply.
_CONTEXT_WARN_RATIO = 0.85

# Base delay (seconds) for the exponential back-off between skill retries.
# Actual delay = _RETRY_BASE_DELAY_S × attempt_number (1-indexed).
_RETRY_BASE_DELAY_S = 0.5

# Capability intents should be answered from the live registry to avoid
# hallucinated tool names and verbose drift.
_TOOLS_QUERY_RE = re.compile(
    r"(what\s+tools|which\s+tools|list\s+available\s+tools|available\s+tools|"
    r"what\s+can\s+you\s+do|capabilities)",
    re.IGNORECASE,
)

# Queries likely to benefit from autonomous multi-source web research.
_RESEARCH_QUERY_RE = re.compile(
    r"(latest|today|current|news|breaking|update|new\s+developments|"
    r"research|analy[sz]e|compare\s+sources|what\s+happened)",
    re.IGNORECASE,
)

_URL_RE = re.compile(r"https?://[^\s)]+", re.IGNORECASE)


def _fmt_error(exc: BaseException) -> str:
    """Return a non-empty error string for UI-safe reporting."""
    msg = str(exc).strip()
    if msg:
        return msg
    return f"{type(exc).__name__}: {exc!r}"


def _truncate_text(text: str, limit: int) -> str:
    if len(text) <= limit:
        return text
    return text[:limit] + f"\n... [truncated at {limit} chars]"


def _read_attachment_text(path: Path, max_chars: int = 12000) -> str:
    """Read an attached file into prompt-safe text."""
    if not path.exists():
        from utils.logger import build_logger
        log = build_logger("agentic.cortex")
        log.warning(f"Attachment file missing: {path}")
        return f"(Attachment not found: {path.name})"
    if not path.is_file():
        from utils.logger import build_logger
        log = build_logger("agentic.cortex")
        log.warning(f"Attachment is not a regular file: {path}")
        return f"(Not a file: {path.name})"

    suffix = path.suffix.lower()
    try:
        file_size = path.stat().st_size
        from utils.logger import build_logger
        log = build_logger("agentic.cortex")
        log.debug(f"Reading attachment: {path.name} ({file_size} bytes, type: {suffix})")
        if suffix in {".txt", ".md", ".py", ".json", ".csv", ".yml", ".yaml", ".toml", ".ini", ".log"}:
            return _truncate_text(path.read_text(encoding="utf-8", errors="replace"), max_chars)

        if suffix == ".pdf":
            from pypdf import PdfReader

            reader = PdfReader(str(path))
            parts: list[str] = []
            for index, page in enumerate(reader.pages, start=1):
                page_text = page.extract_text() or ""
                if page_text.strip():
                    parts.append(f"--- Page {index} ---\n{page_text.strip()}")
            return _truncate_text("\n\n".join(parts) or "(No extractable text found)", max_chars)

        if suffix == ".docx":
            from docx import Document

            document = Document(str(path))
            parts = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
            table_parts: list[str] = []
            for table in document.tables:
                rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
                if rows:
                    header = rows[0]
                    markdown_rows = ["| " + " | ".join(header) + " |", "| " + " | ".join(["---"] * len(header)) + " |"]
                    for row in rows[1:]:
                        markdown_rows.append("| " + " | ".join(row) + " |")
                    table_parts.append("\n".join(markdown_rows))
            return _truncate_text("\n\n".join(parts + table_parts) or "(No extractable text found)", max_chars)

        if suffix == ".xlsx":
            import openpyxl

            workbook = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
            parts: list[str] = []
            for sheet_name in workbook.sheetnames:
                worksheet = workbook[sheet_name]
                rows = []
                for row in worksheet.iter_rows(values_only=True, max_row=50):
                    rows.append([str(cell) if cell is not None else "" for cell in row])
                if not rows:
                    continue
                header = rows[0]
                markdown_rows = ["| " + " | ".join(header) + " |", "| " + " | ".join(["---"] * len(header)) + " |"]
                for row in rows[1:]:
                    markdown_rows.append("| " + " | ".join(row) + " |")
                parts.append(f"### Sheet: {sheet_name}\n" + "\n".join(markdown_rows))
            workbook.close()
            return _truncate_text("\n\n".join(parts) or "(No extractable text found)", max_chars)

        if suffix == ".pptx":
            from pptx import Presentation

            presentation = Presentation(str(path))
            parts: list[str] = []
            for slide_index, slide in enumerate(presentation.slides, start=1):
                slide_text = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
                if slide_text:
                    parts.append(f"--- Slide {slide_index} ---\n" + "\n".join(slide_text))
            return _truncate_text("\n\n".join(parts) or "(No extractable text found)", max_chars)

        return _truncate_text(path.read_text(encoding="utf-8", errors="replace"), max_chars)
    except Exception as exc:
        return f"(Could not extract text from {path.name}: {exc})"


def _build_attachment_context(attachments: list[str]) -> str:
    if not attachments:
        return ""

    import logging
    log = logging.getLogger("agentic.cortex")
    log.info(f"Building attachment context for {len(attachments)} files")
    
    parts: list[str] = ["== Attached Files =="]
    for raw_path in attachments:
        path = Path(raw_path).expanduser().resolve()
        if not path.exists():
            log.warning(f"Attachment file not found: {path}")
        else:
            log.info(f"Including attachment: {path.name}")
        parts.append(f"[File: {path.name}]")
        parts.append(f"Path: {path}")
        parts.append(_read_attachment_text(path))
        parts.append("")
    result = "\n".join(parts).strip()
    log.info(f"Attachment context built: {len(result)} chars")
    return result


class Cortex:
    """
    The reasoning core of Agentic.

    One Cortex instance per application session.  Its async loop runs in
    a background thread and communicates with the UI only via signals.
    """

    def __init__(
        self,
        memory: MemoryLattice,
        registry: SkillRegistry,
        on_token: Callable[[str], None] | None = None,
    ) -> None:
        self._memory   = memory
        self._registry = registry
        self._fabric   = TaskFabric(
            max_concurrent=cfg.get("max_parallel_tasks", 4)
        )
        self._nexus    = GemmaNexus()
        self._weaver   = PromptWeaver(registry.tools_manifest())

        self._loop: asyncio.AbstractEventLoop = asyncio.new_event_loop()
        self._thread = threading.Thread(
            target=self._run_loop, daemon=True, name="cortex-loop"
        )
        self._on_token = on_token   # optional direct callback for UI streaming
        self._active   = False
        # Input queue: replaces the old _deliberating bool so that inputs
        # submitted while deliberation is in progress are queued rather than
        # silently dropped.  _consume_inputs() drains the queue one item at a
        # time, ensuring deliberations are strictly serialised.
        # The queue accepts ``tuple[str, list[str]] | None``; a ``None`` sentinel is pushed by
        # ``stop()`` to wake the consumer and trigger a clean exit without
        # relying on a 1-second polling timeout.
        self._input_queue: asyncio.Queue[tuple[str, list[str]] | None] = asyncio.Queue()
        self._cancel_current = threading.Event()

    async def _stream_tokens(
        self,
        messages: list[dict],
        system_prompt: str,
        temperature: float,
        max_tokens: int,
        tools_schema: list[dict] | None,
    ):
        """
        Stream tokens from ModelNexus with compatibility fallback.

        Some runtime bundles may still have an older ModelNexus.stream()
        signature without ``cancel_event``. In that case, retry without the
        argument so deliberation continues instead of crashing.
        """
        try:
            async for token in self._nexus.stream(
                messages,
                system=system_prompt,
                temperature=temperature,
                max_tokens=max_tokens,
                tools_schema=tools_schema,
                cancel_event=self._cancel_current,
            ):
                yield token
        except TypeError as exc:
            msg = _fmt_error(exc)
            if "unexpected keyword argument 'cancel_event'" not in msg:
                raise
            log.warning(
                "ModelNexus.stream() has no cancel_event support in this runtime; "
                "falling back to legacy stream signature."
            )
            async for token in self._nexus.stream(
                messages,
                system=system_prompt,
                temperature=temperature,
                max_tokens=max_tokens,
                tools_schema=tools_schema,
            ):
                yield token

    def _is_research_query(self, text: str) -> bool:
        return bool(_RESEARCH_QUERY_RE.search(text))

    @staticmethod
    def _extract_urls(text: str) -> list[str]:
        urls = _URL_RE.findall(text)
        deduped: list[str] = []
        seen: set[str] = set()
        for url in urls:
            cleaned = url.rstrip('.,;]')
            if cleaned not in seen:
                seen.add(cleaned)
                deduped.append(cleaned)
        return deduped

    async def _build_auto_research_context(self, query: str) -> str:
        """
        Proactively gather web evidence for research/news requests.

        This stage runs before the regular ReAct loop so the model starts
        with concrete source material instead of only link titles.
        """
        if not cfg.get("deep_research_enabled", True):
            return ""
        if not self._is_research_query(query):
            return ""

        max_results = int(cfg.get("deep_research_max_results", 8))
        max_sources = int(cfg.get("deep_research_max_sources", 4))
        fetch_chars = int(cfg.get("deep_research_fetch_chars", 5000))
        total_chars = int(cfg.get("deep_research_total_chars", 18000))

        lattice.emit_kind(
            SigKind.NOTIFICATION,
            {"message": "Deep research: gathering sources…"},
            source="cortex",
        )

        search_result = await self._registry.invoke(
            "search_web", query=query, max_results=max_results
        )
        if not search_result.success:
            log.warning("Deep research search failed: %s", search_result.error)
            return ""

        candidate_urls = self._extract_urls(str(search_result.output))
        if not candidate_urls:
            return ""

        # Fetch from a larger candidate pool so a few blocked domains (403/429)
        # do not starve the model of source context.
        fetch_pool = candidate_urls[: max(max_sources * 2, max_sources)]

        async def _fetch(url: str) -> tuple[str, str, bool]:
            result = await self._registry.invoke("fetch_web", url=url, max_chars=fetch_chars)
            if result.success:
                return url, str(result.output), True
            return url, f"ERROR: {result.error}", False

        fetched = await asyncio.gather(*[_fetch(url) for url in fetch_pool])

        parts: list[str] = [
            "== Auto Research Brief ==",
            f"Query: {query}",
            f"Sources scanned: {len(fetch_pool)}",
            "",
        ]
        used_chars = sum(len(p) for p in parts)
        success_count = 0
        for idx, (url, text, ok) in enumerate(fetched, start=1):
            block = [
                f"[Source {idx}] {url}",
                "Status: ok" if ok else "Status: error",
                _truncate_text(text, fetch_chars),
                "",
            ]
            block_text = "\n".join(block)
            if used_chars + len(block_text) > total_chars:
                break
            parts.append(block_text)
            used_chars += len(block_text)
            if ok:
                success_count += 1
            if success_count >= max_sources:
                break

        if len(parts) <= 4:
            return ""

        lattice.emit_kind(
            SigKind.NOTIFICATION,
            {"message": f"Deep research: analyzed {len(fetch_pool)} source(s)."},
            source="cortex",
        )
        return "\n".join(parts).strip()

    # ------------------------------------------------------------------
    # Lifecycle
    # ------------------------------------------------------------------

    def start(self) -> None:
        self._active = True
        self._thread.start()
        lattice.attach_loop(self._loop)
        # Schedule the serialised input consumer on the cortex event loop.
        self._loop.call_soon_threadsafe(
            lambda: asyncio.ensure_future(self._consume_inputs())
        )
        log.info("Cortex started (thread: %s)", self._thread.name)

    def stop(self) -> None:
        self._active = False
        # Push a None sentinel to wake _consume_inputs immediately so it can
        # exit cleanly without waiting up to 1 second for a timeout.
        self._loop.call_soon_threadsafe(self._input_queue.put_nowait, None)
        self._loop.call_soon_threadsafe(self._loop.stop)
        self._thread.join(timeout=5)
        log.info("Cortex stopped")

    def update_memory(self, memory: MemoryLattice) -> None:
        """Replace the active MemoryLattice (e.g. when a new session is started)."""
        self._memory = memory

    def _run_loop(self) -> None:
        asyncio.set_event_loop(self._loop)
        self._loop.run_forever()

    # ------------------------------------------------------------------
    # Public entry point (called from UI thread)
    # ------------------------------------------------------------------

    def submit_input(self, text: str, attachments: list[str] | None = None) -> None:
        """
        Submit user input for processing.

        Puts the text onto the input queue.  If a deliberation is already in
        progress the message is held until it completes (no input is dropped).
        When more than one message is already waiting, the user is notified via
        the signal lattice so they know their message is queued.
        """
        if not text.strip() and not attachments:
            return
        
        # Log attachment submission
        if attachments:
            log.info(f"submit_input received {len(attachments)} attachments: {[Path(a).name for a in attachments]}")
        
        payload = (text, attachments[:] if attachments else [])
        queue_size = self._input_queue.qsize()
        if queue_size > 0:
            lattice.emit_kind(
                SigKind.NOTIFICATION,
                {"message": f"Message queued (position {queue_size + 1})."},
                source="cortex",
            )
        # put_nowait is safe: the queue is unbounded and this is a fast O(1) op.
        self._loop.call_soon_threadsafe(self._input_queue.put_nowait, payload)

    def cancel_current(self) -> None:
        """Request cancellation for the currently running model generation."""
        self._cancel_current.set()

    # ------------------------------------------------------------------
    # Deliberation Pulse
    # ------------------------------------------------------------------

    async def _consume_inputs(self) -> None:
        """
        Serialised input consumer – runs on the Cortex event loop.

        Waits for items from *_input_queue* and runs one deliberation at a
        time.  Because this coroutine awaits ``_deliberate_inner`` before
        calling ``get()`` again, deliberations are strictly sequential and
        no input is ever dropped.  The loop exits when a ``None`` sentinel
        is received (pushed by ``stop()``).
        """
        while True:
            item = await self._input_queue.get()
            if item is None:
                # Sentinel: graceful shutdown requested.
                log.debug("Cortex input consumer received shutdown sentinel")
                break
            try:
                user_text, attachments = item
                await self._deliberate_inner(user_text, attachments)
            except Exception as exc:
                log.exception("Unhandled deliberation error: %s", exc)

    async def _deliberate_inner(self, user_text: str, attachments: list[str] | None = None) -> None:
        """Run one end-to-end deliberation pulse for a queued user input."""
        log.info("Deliberation pulse: '%s'", user_text[:80])
        self._cancel_current.clear()

        lattice.emit_kind(
            SigKind.DELIBERATION_START,
            {"input": user_text[:200]},
            source="cortex",
        )

        attachment_block = _build_attachment_context(attachments or [])
        memory_text = user_text
        if attachment_block:
            memory_text = f"{user_text}\n\n{attachment_block}" if user_text.strip() else attachment_block
        self._memory.fluid_write("user", memory_text)

        if self._is_tools_query(user_text):
            final_response = self._build_tools_overview()
            if self._on_token:
                try:
                    self._on_token(final_response)
                except Exception:
                    pass
            self._memory.fluid_write("assistant", final_response)
            lattice.emit_kind(
                SigKind.DELIBERATION_END,
                {"response": final_response, "skills_used": 0},
                source="cortex",
            )
            return

        mem_ctx = self._memory.assemble_context(
            include_crystal=5, include_bedrock=10, query=user_text
        )
        system_prompt = self._weaver.build_system(memory_context=mem_ctx)
        fluid = self._memory.fluid_read()
        messages = self._weaver.build_messages(fluid[:-1], user_text)

        if attachment_block:
            messages.insert(len(messages) - 1, {"role": "user", "content": attachment_block})

        research_block = await self._build_auto_research_context(user_text)
        if research_block:
            messages.insert(len(messages) - 1, {"role": "user", "content": research_block})

        specs = self._registry.all_specs()
        autopilot_enabled = bool(cfg.get("autopilot_enabled", True))
        current_model = cfg.get("model_id", "google/gemma-3-1b-it")
        plan = build_plan(
            user_text,
            specs,
            current_model=current_model,
            has_attachments=bool(attachments),
        )

        if autopilot_enabled:
            route_block = format_plan_block(plan)
            if route_block:
                messages.insert(len(messages) - 1, {"role": "user", "content": route_block})

        tools: list[dict] | None = (
            self._registry.tools_schema()
            if self._nexus.tool_calls_supported
            else None
        )

        context_limit: int = cfg.get("context_limit_tokens", 4096)
        estimated_tokens = self._estimate_tokens(messages, system_prompt)
        if estimated_tokens > context_limit * _CONTEXT_WARN_RATIO:
            before_tokens = estimated_tokens
            before_count = len(messages)
            estimated_tokens = self._trim_messages_to_context(
                messages, system_prompt, context_limit
            )
            trimmed = before_count - len(messages)
            log.warning(
                "Prompt size (%d tokens) exceeded %.0f%% of context limit (%d). "
                "Trimmed %d oldest message(s) → ~%d tokens remain.",
                before_tokens,
                _CONTEXT_WARN_RATIO * 100,
                context_limit,
                trimmed,
                estimated_tokens,
            )

        max_iterations = cfg.get("react_max_iterations", 6)
        if autopilot_enabled and plan.long_horizon:
            max_iterations = min(20, max_iterations + 2)
        generation_temperature = float(cfg.get("generation_temperature", 0.25))
        generation_max_tokens = int(cfg.get("generation_max_tokens", 512))

        final_response = ""
        total_skills_used = 0
        total_skill_failures = 0

        model_candidates = plan.model_candidates if autopilot_enabled else [current_model]
        model_idx = 0
        current_route_model = model_candidates[model_idx]
        checkpoint_every = int(cfg.get("autopilot_checkpoint_every_n", 1))
        escalate_enabled = bool(cfg.get("autopilot_escalation_enabled", True))
        escalate_ratio = float(cfg.get("autopilot_escalate_on_error_ratio", 50)) / 100.0

        self._nexus.set_model_override(current_route_model)

        try:
            for iteration in range(max_iterations):
                response_parts: list[str] = []

                try:
                    async for token in self._stream_tokens(
                        messages,
                        system_prompt,
                        generation_temperature,
                        generation_max_tokens,
                        tools,
                    ):
                        response_parts.append(token)
                        if self._on_token:
                            try:
                                self._on_token(token)
                            except Exception:
                                pass
                        if self._cancel_current.is_set():
                            break
                except Exception as exc:
                    err_text = _fmt_error(exc)
                    log.exception("Model stream failed: %s", err_text)
                    err_msg = f"\n[Model error: {err_text}]"
                    response_parts.append(err_msg)
                    if self._on_token:
                        self._on_token(err_msg)

                full_response = "".join(response_parts)
                if self._cancel_current.is_set():
                    final_response = full_response if full_response else "Response stopped by user."
                    break

                final_response = full_response
                skill_queue = self._weaver.extract_skill_calls(full_response)
                if not skill_queue:
                    break

                results = await self._run_skills(skill_queue)
                total_skills_used += len(skill_queue)
                failures = sum(1 for _, _, ok in results if not ok)
                total_skill_failures += failures

                lattice.emit_kind(
                    SigKind.REACT_ITERATION,
                    {
                        "iteration": iteration + 1,
                        "skills_run": [inv.skill_name for inv in skill_queue],
                        "result_count": len(results),
                        "model": current_route_model,
                    },
                    source="cortex",
                )

                if autopilot_enabled and ((iteration + 1) % checkpoint_every == 0):
                    cp_text = (
                        f"task={plan.task_kind}; iter={iteration + 1}; "
                        f"model={current_route_model}; skills={len(skill_queue)}; "
                        f"failures={failures}"
                    )
                    self._memory.bedrock_write("autopilot_checkpoint", cp_text, confidence=0.55)

                if self._on_token and iteration + 1 < max_iterations:
                    skill_names = ", ".join(inv.skill_name for inv in skill_queue)
                    try:
                        self._on_token(_REACT_ITERATION_STATUS.format(skill_names=skill_names))
                    except Exception:
                        pass

                if (
                    autopilot_enabled
                    and escalate_enabled
                    and len(skill_queue) > 0
                    and (failures / len(skill_queue)) >= escalate_ratio
                    and (model_idx + 1) < len(model_candidates)
                ):
                    model_idx += 1
                    current_route_model = model_candidates[model_idx]
                    self._nexus.set_model_override(current_route_model)

                asst_role = get_assistant_role(current_route_model)
                messages.append({"role": asst_role, "content": full_response})
                obs_text = self._weaver.format_observations(results)
                messages.append({"role": "user", "content": obs_text})
            else:
                asst_role = get_assistant_role(current_route_model)
                messages.append({"role": asst_role, "content": final_response})
                messages.append({"role": "user", "content": _REFLEXION_PROMPT})
                if self._on_token:
                    try:
                        self._on_token(_REFLEXION_STATUS)
                    except Exception:
                        pass
                reflect_parts: list[str] = []
                try:
                    async for token in self._stream_tokens(
                        messages,
                        system_prompt,
                        generation_temperature,
                        generation_max_tokens,
                        None,
                    ):
                        reflect_parts.append(token)
                        if self._on_token:
                            try:
                                self._on_token(token)
                            except Exception:
                                pass
                        if self._cancel_current.is_set():
                            break
                except Exception as exc:
                    log.exception("Reflexion stream failed: %s", exc)
                    reflect_parts.append(f"\n[Reflexion error: {exc}]")
                if reflect_parts:
                    final_response = "".join(reflect_parts)

            if (
                autopilot_enabled
                and bool(cfg.get("autopilot_quality_gate_enabled", True))
                and (model_idx + 1) < len(model_candidates)
            ):
                gate = quality_gate(
                    final_response,
                    query=user_text,
                    tool_calls=total_skills_used,
                    skill_failures=total_skill_failures,
                )
                if not gate.passed:
                    model_idx += 1
                    current_route_model = model_candidates[model_idx]
                    self._nexus.set_model_override(current_route_model)
                    asst_role = get_assistant_role(current_route_model)
                    messages.append({"role": asst_role, "content": final_response})
                    messages.append(
                        {
                            "role": "user",
                            "content": (
                                "Quality gate requested a cleaner final response. "
                                "Provide a direct, complete answer now with no tool calls."
                            ),
                        }
                    )
                    upgrade_parts: list[str] = []
                    try:
                        async for token in self._stream_tokens(
                            messages,
                            system_prompt,
                            generation_temperature,
                            generation_max_tokens,
                            None,
                        ):
                            upgrade_parts.append(token)
                            if self._on_token:
                                try:
                                    self._on_token(token)
                                except Exception:
                                    pass
                    except Exception as exc:
                        log.exception("Quality fallback stream failed: %s", exc)
                    if upgrade_parts:
                        final_response = "".join(upgrade_parts)

            if not final_response:
                final_response = (
                    "Error: the model did not generate a response. "
                    "Please try again or rephrase your input."
                )

            self._memory.fluid_write("assistant", final_response)
            lattice.emit_kind(
                SigKind.DELIBERATION_END,
                {
                    "response": final_response,
                    "skills_used": total_skills_used,
                    "skill_failures": total_skill_failures,
                    "model": current_route_model,
                    "task_kind": plan.task_kind,
                },
                source="cortex",
            )
        finally:
            self._nexus.set_model_override(None)

    def _is_tools_query(self, text: str) -> bool:
        return bool(_TOOLS_QUERY_RE.search(text))

    def _build_tools_overview(self) -> str:
        specs = self._registry.all_specs()
        if not specs:
            return "No tools are currently registered."

        lines = [
            "Here are the available tools and when I use each:",
            "",
        ]
        for spec in specs:
            desc = spec.description.strip().rstrip(".")
            lines.append(f"- {spec.name}: {desc}.")
        lines.append("")
        lines.append("I use these automatically when your request clearly needs them.")
        return "\n".join(lines)

    # ------------------------------------------------------------------
    # Context-window helpers
    # ------------------------------------------------------------------

    def _estimate_tokens(self, messages: list[dict], system_prompt: str) -> int:
        """
        Count prompt tokens using the loaded tokenizer (accurate), falling
        back to the chars-per-token heuristic when the model is not yet loaded.
        """
        return (
            sum(self._nexus.count_tokens(m.get("content", "")) for m in messages)
            + self._nexus.count_tokens(system_prompt)
        )

    def _trim_messages_to_context(
        self,
        messages: list[dict],
        system_prompt: str,
        context_limit: int,
    ) -> int:
        """
        Trim *messages* in-place from the oldest end until the estimated
        token count is within *context_limit*.

        The current user turn (last element) is always preserved so that the
        model always sees the latest request.  Returns the final token count
        after trimming.
        """
        while len(messages) > 1:
            tokens = (
                sum(self._nexus.count_tokens(m.get("content", "")) for m in messages)
                + self._nexus.count_tokens(system_prompt)
            )
            if tokens <= context_limit:
                break
            messages.pop(0)
        return (
            sum(self._nexus.count_tokens(m.get("content", "")) for m in messages)
            + self._nexus.count_tokens(system_prompt)
        )

    # ------------------------------------------------------------------
    # Skill execution primitive
    # ------------------------------------------------------------------

    async def _run_skills(
        self,
        invocations: list[SkillInvocation],
    ) -> list[tuple[SkillInvocation, str, bool]]:
        """
        Execute skill calls as parallel Task Fibers and return structured results.

        Each invocation gets its own fiber.  All fibers are enqueued and then
        run_until_empty() drives them to completion concurrently (up to
        max_concurrent from the TaskFabric configuration).

        Returns a list of ``(invocation, result_str, success)`` tuples in the
        same order as *invocations*.  This is consumed by
        ``PromptWeaver.format_observations`` to build the Observation turn for
        the next ReAct iteration.
        """
        id_to_inv: dict[str, SkillInvocation] = {}
        for inv in invocations:
            fiber = TaskFiber(
                label=f"skill:{inv.skill_name}",
                fn=self._make_skill_fn(inv),
                tags=["skill"],
            )
            fid = self._fabric.add_fiber(fiber)
            id_to_inv[fid] = inv

        await self._fabric.run_until_empty()

        results: list[tuple[SkillInvocation, str, bool]] = []
        for fid, inv in id_to_inv.items():
            fiber = self._fabric.get_fiber(fid)
            if fiber is None:
                results.append((inv, "ERROR: fiber not found", False))
                continue
            if fiber.status == FiberStatus.DONE:
                results.append((inv, str(fiber.result)[:1000], True))
            elif fiber.status == FiberStatus.FAILED:
                results.append((inv, f"ERROR: {fiber.error}", False))
            else:
                results.append(
                    (inv, f"ERROR: unexpected fiber state {fiber.status}", False)
                )
        return results

    def _make_skill_fn(self, inv: SkillInvocation):
        async def _fn(fiber: TaskFiber):
            retry_budget: int = cfg.get("skill_retry_budget", 1)
            fiber.set_progress(0.1)
            last_result = None
            for attempt in range(retry_budget + 1):
                last_result = await self._registry.invoke(inv.skill_name, **inv.args)
                if last_result.success:
                    fiber.set_progress(1.0)
                    return last_result.output
                if attempt < retry_budget:
                    delay = _RETRY_BASE_DELAY_S * (attempt + 1)
                    log.warning(
                        "Skill '%s' failed (attempt %d/%d): %s — retrying in %.1fs…",
                        inv.skill_name, attempt + 1, retry_budget + 1,
                        last_result.error, delay,
                    )
                    await asyncio.sleep(delay)
            raise RuntimeError(
                last_result.error if last_result else "skill returned no result"
            )
        return _fn
